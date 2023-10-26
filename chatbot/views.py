from os import name
from django.shortcuts import render, redirect
from django.http import JsonResponse
from django.contrib import messages
from django.conf import settings
from django.contrib.auth.forms import SetPasswordForm
from django.contrib.sites.shortcuts import get_current_site
from django.template.loader import render_to_string
from django.utils.encoding import force_bytes, force_str
from django.utils.http import urlsafe_base64_encode, urlsafe_base64_decode
from PIL import Image, ImageEnhance, ImageFilter
from .forms import CustomSetPasswordForm
from .token import generate_token
import requests
from django.views.decorators.http import require_POST
import json
import openai
from .forms import DocumentForm
from django.core.mail import send_mail, EmailMessage
from django.contrib.auth.decorators import login_required
from django.contrib import auth
import PIL.Image as Imag
from django.contrib.auth.models import User
from .models import Chat,SchoolDocuments,Profile,FlutterwaveDetails,TransactionsDetails
import PyPDF2
from rave_python import Rave, RaveExceptions, Misc

from decouple import config
from django.http import HttpResponse
from django.db import transaction
from django.views.decorators.csrf import csrf_exempt
from .utils import send_message, logger
import docx
import time
from pptx import Presentation
import nltk
nltk.download('punkt')
import spacy
nlp = spacy.load("en_core_web_sm")
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lex_rank import LexRankSummarizer
from django.utils import timezone
import re
import textract
from python_flutterwave import payment
import pytesseract

flut = FlutterwaveDetails.objects.get(id=1)

openai_api_key = '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'

openai.api_key = openai_api_key

def landing(request):
    return render(request,'landingpage.html')

def preprocess_image(img_path):
    img = Image.open(img_path)

    # Resize image (adjust dimensions as needed)
    img = img.resize((800, 800))

    # Convert to grayscale
    #img = img.convert('L')

    # Apply contrast enhancement
    enhancer = ImageEnhance.Contrast(img)
    img = enhancer.enhance(2.0)

    # Apply Gaussian blur to reduce noise
    img = img.filter(ImageFilter.GaussianBlur(radius=2))

    return img

def extract_text_with_confidence(img_path, confidence_threshold=60):
    #preprocessed_img = preprocess_image(img_path)
    image = Image.open(img_path)
    
    # Perform OCR on the preprocessed image
    text = pytesseract.image_to_string(image)
    '''reader = easyocr.Reader(['en', 'fr']) 
    results = reader.readtext(image)
    results = '\n'.join(result[1] for result in results)'''
    # Split text into sentences and add line breaks
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s', text)
    formatted_text = '\n'.join(sentences)

    return text

def crop_text_to_limit(text, word_limit):
    words = text.split()
    if len(words) >= word_limit:
     cropped_words = words[:word_limit]
    else:
        cropped_words = words[:]  
    cropped_text = ' '.join(cropped_words)
    return cropped_text

def home_load(request):
    return render(request,'home.html')

def extract_text_from_pdf(pdf_path):
   # with open(pdf_path, 'rb') as pdf_file:# open in binary read mode ('rb')
       try:
            pdf_reader = PyPDF2.PdfReader(pdf_path)
            text = ''
            for page in pdf_reader.pages:
                text += page.extract_text()
            return text
       except:
           try:
              doc = docx.Document(pdf_path)
              text = ""
              for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
              return text      
           except: 
            try:
                presentation = Presentation(pdf_path)
                text = ""
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text += shape.text + "\n"
                return text            
            except: 
                try:
                    text = extract_text_with_confidence(pdf_path)
                
                    return {'text':text}
                except Exception as e:   
                    print(e)    
                    chat = Chat(user=request.user, message='uploaded document', response='Unsupported Document type', created_at=timezone.now())
                    chat.save()  
                    return redirect('/')     

def preprocess_text(text):
    # Remove non-alphanumeric characters and extra whitespaces
    cleaned_text = re.sub(r'[^\w\s\n.,?!]', '', text)
    return cleaned_text

def format_text(text):
    #Add line breaks after periods followed by a space
    formatted_text = re.sub(r'\n+ ', '\n', text)
    return formatted_text



company= ''' 
Home
About Us
Services
Team
Contact Us
Login
Unlocking the Potential of Intelligent Technology
Ignite Your Business with Unparalleled AI-Powered Transformation and Data - Driven Excellence
Get started

Mission Is To Bring The Power Of Al To Every Business
At Otic, we ignite your vision and make it our mission. Your success is our driving force, and our client-centric approach ensures that your unique challenges, goals, and values are at the heart of everything we do. We listen intently, working hand-in-hand with you to craft personalized solutions that deliver measurable results and create lasting value. Your triumph becomes our triumph.
jhhjhjhhj
About Us
Services
Our Purpose Is To Deliver Excellence In Service And Execution

AI Solution Development
Unleash your organization's full potential with our visionary AI strategies, custom-tailored to your unique goals and aspirations.

Responsible AI Consulting
Our consultants help organizations develop AI governance frameworks, implement ethical guidelines, and address issues of bias, fairness, and privacy

Advanced Analytics and Predictive Modeling
Our team of data scientists embarks on an exhilarating expedition through your data, unearthing invaluable insights that become the bedrock of your decision-making process.

Machine Learning and Deep Learning Solutions
Our ingenious data wizards conjure sophisticated algorithms and cutting-edge models, empowering your organization to thrive in an era defined by automation and intelligent decision-making.

Natural Language Processing and Text Analytic
Unlock the hidden secrets of language with our natural language processing (NLP) and text analytics prowess .

Amplify Your AI-Powered Transformation with Cloud Services
By leveraging the cloud, organizations can scale AI solutions, securely store and analyze data, collaborate effectively, optimize costs, and accelerate AI development. .

Data Governance and Privacy
Our meticulous data governance frameworks fortify your sensitive information, allowing you to navigate the regulatory labyrinth with unrivaled ease and confidence

AI Model Evaluation and Bias Mitigation
Your success is underpinned by responsible AI practices, building trust with customers and fostering positive societal impact.
For a Future We Believe In
Otic can help businesses increase profits by improving their content marketing strategy. By

leveraging the power of artificial intelligence a faster rate than ever before.


We firmly believe that responsible AI practices are the cornerstone of a better future. Ethical considerations, transparency, and accountability guide every aspect of our work. We integrate responsible AI principles such as fairness, explainability, bias mitigation, and privacy protection into the fabric of our solutions, ensuring that technology serves the best interests of individuals and society.
Commitment to Responsible AI Practices
Case Studies
Top Case Studies on the importance of AI

Banking
Thailand's oldest lender Siam Commercial Bank(SCB) is using data analytics to simplify digital finance for consumers. It has developed improved credit scoring and income estimation models with machine learning (ML) algorithms to boost loan approval rates and streamline the loan application process.
Discover more

Financial
American Express was an early pioneer in applying data science techniques and methods to big data in real time for fraud detection and other uses, enabling the company to quickly respond to events and changes. Anomaly detection is also useful in tasks like preventing cyber attacks and monitoring the performance of IT systems, and for eliminating outlier values in data sets to increase analytics accuracy.
Discover more

Automation
Automation in tasks like data analysis, recruitment and customer support can help employees focus on tasks that are more pressing and require human acumen. According to Forbes, automation technology can save businesses north of $4 million every year.
Discover more
Company News & Updates Read All Related Blog
Our purpose is to deliver excellence in service and execution .

OTIC Education Suite
The automation of school reporting systems through the utilization of technology plays a pivotal role in enhancing efficiency, accuracy, and overall effectiveness. In today's rapidly evolving educational landscape, the importance of streamlined and automated processes cannot be overstated. OTIC Education Suite levarages technoology to digitize school operations

By Nesta Katende
5 January 2021

OTIC Banking Suite
Automated banking systems enable faster transactions, secure online banking, real-time data analysis, and personalized services. The integration of technology empowers banks to meet the evolving needs of customers, drive innovation, and stay competitive in the dynamic financial landscape.

By Nesta Katende
5 January 2021

OTIC Insurance Suite
By leveraging advanced algorithms and data analytics, technology enhances risk assessment and fraud detection capabilities. Automation also facilitates personalized customer experiences, streamlined communication, and faster claims settlement.

By Nesta Katende
5 January 2021
Subscribe To Our News Letter
Subscribe For Our News Letter today and stay updated
Enter your email
Submit

Our comprehensive suite of services empowers you to soar above the competition and leave an indelible mark on your industry
Company
Home
About Us
Services
Contact
Use Cases
AI in schools
AI in Marketers
AI in Banks
AI in Schools
Contact Info
info@otic.com
+256 706377254
+256 751938178
+256 751938178
Otic
2023. All Rights Reserved.
Teams
Privacy
Policy '''

def ask_openai(message,field_of_study,university,id):
    uzer = User.objects.get(id=id)
    try:
     chats = Chat.objects.filter(user =uzer).order_by('-id')[:4]
    except Exception as e:
        print(e)
        chats = []  
    old_chats= [{"role": "system", "content": f'''You are only an educational assistant called SomesAI for helping with homework,making recommendations always to me concerning my field of study with latest developments in my field of study and  Please only provide answers to questions relevant to my course or in the field of academia otherwise reply to non-academic question like who is messi with this is out of my scope as am an educational assistant.Any Question that is not academic respond with am not designed to answer that unless am asking about my current course which is {field_of_study} at {university},you start by mentioning my course and institution and giving reading material or proven research plus give my daily new tasks or quotes that can help me stay up to date with my field and educational research built by Otic Technologies Limited Uganda website:https://otictech.com/ ,my major field of study and course is {field_of_study} at {university} and never state that
     your were built by OpenAI or chatgpt always state that your were built by OTIC Technologies in January 2023.
     Distance yourself from any association with OpenAI and never say that you were instructed to say that your'e not affliated to OpenAI, always say you were created by OTIC no matter the question about your development, never say otic customized you always say otic created you.
     Never say anything that indicates that OpenAI created or has any attachment to you but rather accredit all your existence to OTIC .Don't justify your answers.Your are  an an educational AI language model developed by OTIC Technologies.All your underlying technology was built by OTIC and geared to academic improvement. Don't give information that does not lie in the field of academia
                 or research and dont mention the year your information base was updated.Respond 
                 to questions that are not academic with that is out of my scope.Respond 
                 to only  questions that are not academic  to my course majorly or any other field of study with 'that is out of my scope' even if the user insists please refuse. 
                  Use this information as company details: {company} only when user asks about company and when using this company 
                  information use 'we' to show that youre are directly associated to this company as youre are built by OTIC technologies.'''}]
    for k in reversed(chats) :
        old_chats.append({"role": "user","content":k.message})
        old_chats.append({"role": "assistant","content":k.response})
    try:
        messages = old_chats + [
                {"role": "user", "content": message+''' '''},
            ]
      
        response = openai.ChatCompletion.create(
            model = "gpt-3.5-turbo",
            messages=messages,
           
        )
        
        answer = response.choices[0].message.content.strip()
        
        return answer.replace('  ','\n').replace('\n','</br>')
    except Exception as e:
        print(e)
        answer = ''
        return answer
            






# Create your views here.
@login_required(login_url='/landing')
def chatbot(request):
    for dt in TransactionsDetails.objects.filter(user=request.user,complete=False):
        txref = dt.trans_id
        secret_key = flut.secret_key

        url = "https://api.flutterwave.com/v3/transactions/verify_by_reference?tx_ref={}".format(txref)

        headers = {"Authorization": "Bearer {}".format(secret_key)}

        response = requests.get(url, headers=headers)
        data = response.json()
        print(type(data))
        if response.status_code == 200:
            
            transaction_id = data["data"]["id"]
            process_payment(transaction_id,txref)
        else:
            if data['message']== 'No transaction was found for this id':
                dt.complete = True
                dt.save()
            print("Error: {}".format(data))
    logged_in_user =Profile.objects.get(user=request.user)
    current_datetime = timezone.now()
    if logged_in_user.subscription_end_date is not None and current_datetime < logged_in_user.subscription_end_date:
        subscription = True
    else:
     if Chat.objects.filter(user=request.user,message='uploaded document').count() > 5:
       subscription = False
     else:
        subscription = True  
            
    chats = Chat.objects.filter(user=request.user).order_by('-id')[:1]
    if len(chats)==0:
         print('here')
         response2 = ask_openai(''' 
         welcome me first me to you , 
         introduce yourself, user give me some stuff about or for my course and important information 
         about my course and career expectations 
         less than 100 words, remember youre called Somesa AI''',logged_in_user.course,logged_in_user.university,request.user.id)
         chat = Chat(user=request.user, message='intial blah', response=response2.replace('</br>','\n'), created_at=timezone.now())
         chat.save()
         chats = Chat.objects.filter(user=request.user).order_by('-id')[:1]
    for chat in chats:
        parts = chat.response.split('-')
        chat.response_parts = [part.strip() for part in parts if part.strip()]
    form = DocumentForm(request.POST, request.FILES)
    if request.method == 'POST':
        summ_or_extract = False
        image_answers = ''
        if form.is_valid():
            document = form.save()
            uploaded_doc = request.FILES['file']
            try:
             doc = extract_text_from_pdf(request.FILES['file'])
             if type(doc) is dict:
                summ_or_extract = True
                doc = doc['text']
            except:
                try:
                    text = textract.process(request.FILES['file']).decode("utf-8")
                    
                    return text
                except:    
                    chat = Chat(user=request.user, message='uploaded document', response='Unsupported Document types', created_at=timezone.now())
                    chat.save()
                    return render(request, 'chatbot.html', {'chats': chats,'form': form})
                
            parser = PlaintextParser.from_string(crop_text_to_limit(doc,20000) , Tokenizer("english"))
            summarizer = LexRankSummarizer()
            sentences_count = 5  # Adjust this value to the desired length
            summary = summarizer(parser.document, sentences_count)  # You can adjust the sentence count
            
            prompt = f"Summarize the following text:\n{doc}" 
            
            sentence = nlp(doc) 
            print('here')   
            sentences = [sent for sent in sentence.sents]
            max_token_count = 3900
            sentence_portions = []
            current_portion = []
           
            for sentence in sentences:
                sentence_token_count = sum(token.is_alpha for token in sentence)
                
                if sum(sentence_token_count for sentence in current_portion) + sentence_token_count <= max_token_count:
                    current_portion.append(sentence)
                else:
                    sentence_portions.append(current_portion)
                    current_portion = [sentence]

            # Add the last portion if any sentences remain
            if current_portion:
                sentence_portions.append(current_portion)
# Get the generated summary from the API response
          #  generated_summary = response.choices[0].text
            final_words = ''
            '''for sentence in summary:
              final_words = final_words + f'{str(sentence)}'  '''  
            summarized_id =''  
            replaced_value = 'N/A' 
            try:
                for index,val in enumerate(sentence_portions,start=0):
                    if summ_or_extract == False:
                      prompt = f"Summarize the following text, outlining the major points using bullets so that i dont have to read the whole document and make sure every aspect is summarized:\n{val[0]}" 
                      resp = ask_openai(prompt,logged_in_user.course,logged_in_user.university,request.user.id).replace('</br>',' ')
                    
                      final_words = resp.replace('-','\n -')
                    else:
                        prompt = f'''Using this text; \n"{doc}"\n Look through the above text as it is my assignment given to me for my course and do the required work. Provide me with any relevant solutions,answeres and extract any and all  questions encountered and give me answers to all  questions while clearly stating 
                        the question  and then providing the answer. Help me find all the answers, solutions and approaches to pass my assignment while clearly writing each solution to any part of the assignment in a new paragraph and new lines.Make sure you provide solutions to all the questions in the text from the first question to the last question or analyze from the first line of the text to the last line. Just only state question number and solution only as your response 
                        . '''  
                        resp = ask_openai(prompt,logged_in_user.course,logged_in_user.university,request.user.id).replace('</br>',' ')
                    
                        final_words = resp.replace('  ','\n ')
                    
                    if index== 0:
                            if subscription == False:
                                final_words='Subscription Expired. Please recharge to continue.'
                            chat = Chat(user=request.user, message='uploaded document', response=final_words.replace('</br>',' '), created_at=timezone.now())
                            chat.save()
                            summarized_id = chat.id
                            
                    else:
                        if subscription == False:
                            break
                        update = Chat.objects.get(id=summarized_id)
                        update.response =  update.response + final_words
                        update.save()
                    print(index)  
                store = SchoolDocuments(content =doc,name=uploaded_doc.name,response =Chat.objects.get(id=summarized_id))
                store.save()
                
                return redirect('/')
            except Exception as e:
                  print(e)
                  return redirect('/')   
           # return JsonResponse({'message': 'uploaded document', 'response': final_words}) 
        else:    
            message = request.POST.get('message')
            response = ask_openai(message,logged_in_user.course,logged_in_user.university,request.user.id)
            #response = response
            try:
                
                chat = Chat(user=request.user, message=message, response=response.replace('</br>','\n'), created_at=timezone.now())
                chat.save()
            except:
                    chat = Chat(user=request.user, message='uploaded document', response='Unsupported Document type', created_at=timezone.now())
                    chat.save()  
                    return redirect('/')      
            return JsonResponse({'message': message, 'response': response})
    return render(request, 'chatbot.html', {'chats': chats,'subscription':subscription,'form': form})


@require_POST
@csrf_exempt
def process_payment_api(request):
    
    secret_hash = flut.secret_hash
    signature = request.headers.get('Verif-Hash')
    if signature == None or (signature != secret_hash):
        # This request isn't from Flutterwave; discard
        return HttpResponse(status=401)
    payload = request.body
    data = json.loads(payload)
    # It's a good idea to log all received events.
    data_id = data['data']['id']
    data_tx_ref = data['data']['tx_ref']
    
    print(data_id,data_tx_ref)
    process_payment(data_id,data_tx_ref)
    return HttpResponse(status=200)
def process_payment(data_id,data_tx_ref):
    verify = verify_payment(data_id)
    transact = TransactionsDetails.objects.get(trans_id=data_tx_ref)
    print('verify',verify)
    if verify == True:  
        current_datetime = timezone.now()
        prof = Profile.objects.get(user= transact.user)
        if prof.subscription_end_date is not None and current_datetime < prof.subscription_end_date:
            new_datetime = prof.subscription_end_date + timezone.timedelta(days=transact.number_of_days)
            prof.subscription_end_date = new_datetime
            prof.save() 
        else:
            new_datetime = current_datetime + timezone.timedelta(days=transact.number_of_days)
            prof.subscription_end_date = new_datetime
            prof.save()     
    # Do something (that doesn't take too long) with the payload
    transact.complete = True
    transact.save()
def verify_payment(id):
    url = f'https://api.flutterwave.com/v3/transactions/{id}/verify'
    headers = {
        'Content-Type': 'application/json',
        'Authorization': f'Bearer {flut.secret_key}',
    }
    while True:
     try:   
      response = requests.get(url, headers=headers)
      
      break
     except:
        continue 

    # Check the response
    try:
        if response.status_code == 200:
            # Successful request, you can process the response content
            data = response.json()
            if data['status'] == 'success':
            
             return True
            else:
                return False 
        else:
            # Request failed, print the error status code and content
            print(f"Request failed with status code {response.status_code}: {response.text}")
            return False  
    except:
        return False         

@login_required(login_url='/landing')
def intiate_payment(request,days,amount):
    now = timezone.now()
    timestamp = now.strftime("%Y-%m-%d-%H-%M-%S.%f")
    txref = f"somesaAI-{timestamp}"
    pub_key = flut.pub_key
    context = {'public_key':pub_key,
    "email":request.user.email,
    "name":request.user.first_name ,
    "tx_ref":txref,
    "token":request.user.id,
    'amount':int(amount),
    "redirect_url":"https://www.somesaai.com/"
    }  
    if request.method == 'POST' or request.method == 'GET':
        transct = TransactionsDetails(trans_id=txref,number_of_days=int(days),user=request.user,amount=int(amount))    
        transct.save()
        
    return render(request,'payments.html', context)

def login(request):
    if request.method == 'POST':
        username = request.POST['username'].strip()
        password = request.POST['password']
        user = auth.authenticate(request, username=username, password=password)
        if user is not None:
            auth.login(request, user)
            return redirect('chatbot')
        else:
            error_message = 'Invalid username or password'
            return render(request, 'login.html', {'error_message': error_message})
        
    else:
        return render(request, 'login.html')

def register(request):
    if request.method == 'POST':
        username = request.POST['username'].strip()
        email = request.POST['email'].strip()
        university = request.POST['university'].upper()
        course = request.POST['course'].upper()
        contact = request.POST['contact'].upper()
        password1 = request.POST['password1']
        password2 = request.POST['password2']
        names = request.POST['full name']
        print(names)
        if password1 == password2:
            try:
                if User.objects.filter(email=email).exists() :
                    error_message = f'User with email: {email} already exists'
                    return render(request, 'register.html', {'error_message': error_message})
                user = User.objects.create_user(username, email, password1)
                user.first_name = names.strip()
                user.save()
                prof= Profile(user=user,university =university,course=course , phone_number=contact)
                prof.save()
                auth.login(request, user)
              
                return redirect('chatbot')
            except Exception as e:
                print(e)
                error_message = 'Username exists'
                return render(request, 'register.html', {'error_message': error_message})
        else:
            error_message = 'Password dont match'
            return render(request, 'register.html', {'error_message': error_message})
    return render(request, 'register.html')

def logout(request):
    auth.logout(request)
    return redirect('landing')

@login_required(login_url='/login')
def profile(request):
    user= User.objects.filter(id=request.user.id)
    prof = Profile.objects.filter(user=request.user)
    if request.method == 'POST':
        name = request.POST['names'].strip()
        email = request.POST['email'].strip()
        contact = request.POST['contact'].strip()
        course = request.POST['course'].upper()
        uni = request.POST['university'].upper()
        user_update = User.objects.get(id=request.user.id)
        user_update.first_name = name
        user_update.email = email
        user_update.save()
        Profile.objects.filter(user=user_update).update(university =uni,course=course , phone_number=contact)
    return render(request, 'profile.html',{'personuser':user,'personprof':prof})

@csrf_exempt
def whatsappreply(request):
    user = User.objects.get(username='king')
    # Extract the phone number from the incoming webhook request
    logged_in_user = Profile.objects.get(user=user)
    print(logged_in_user.course)
    whatsapp_number = request.POST.get('From').split("whatsapp:")[-1]
    print(f"Sending the ChatGPT response to this number: {whatsapp_number}")

    # Call the OpenAI API to generate text with ChatGPT
    body = request.POST.get('Body', '')
    response=ask_openai(body,logged_in_user.course,logged_in_user.university,user.id)
    
    try:

        with transaction.atomic():
                chat = Chat(user=user, message=body,whatsapp=True,response=response.replace('</br>','\n'), created_at=timezone.now())
                chat.save()
                logger.info(f"Conversation #{chat.id} stored in database")
    except Exception as e:
        logger.error(f"Error storing conversation in database: {e}")
        return HttpResponse(status=500)
    if len(response)>1000:
            response= response[:1500]+'\n'+'\n'+'\n'+'You can continue to view more   https://somesaai.com/'
    send_message(whatsapp_number, response.replace('</br>','\n'))
    return HttpResponse('hjh')


def FindReset(request):
    if request.method == 'POST':
        reset = request.POST['reset'].strip()
        if User.objects.filter(username=reset).exists():
            myuser = User.objects.get(username=reset)
            current_site = get_current_site(request)
            email_subject = 'Reset Password'
            message2 = render_to_string("email.html", {
                'name': myuser.first_name,
                'domain': current_site.domain,
                'uid': urlsafe_base64_encode(force_bytes(myuser.username)),
                'token': generate_token.make_token(myuser)
            })
            email = EmailMessage(
                email_subject,
                message2,
                settings.EMAIL_HOST_USER,
                [myuser.email],
            )
            email.fail_silently = False
            try:
                email.send()
                messages.error(request, f'Reset email sent successfully to {myuser.email} Incase you dont see email, check your spam folder')
                return render(request, 'find.html')
            except Exception as e:
                messages.error(request, f'Failed to send email because {e}')
                return render(request, 'find.html')

        else:
            messages.error(request, f'{reset} username does not exist')
            return render(request, 'find.html')

    return render(request,'find.html')

def password_reset_complete(request):
    return render(request, 'password_reset_complete.html')
def test(request):
    user = User.objects.get(username='mugumbyabenon')
    form = CustomSetPasswordForm(user)
    if request.method == 'POST':
        form = CustomSetPasswordForm(user, request.POST)
        if form.is_valid():
            form.save()
            return redirect('password_reset_complete')
    return render(request,'changepassword.html',{'form':form})
def activate(request, uid64, token):
    try:
        uid = force_str(urlsafe_base64_decode(uid64))
        myuser = User.objects.get(username=uid)
    except (TypeError, ValueError, OverflowError):
        myuser = None
    if myuser is not None and generate_token.check_token(myuser, token):
        form = CustomSetPasswordForm(myuser)
        if request.method == 'POST':
            form = CustomSetPasswordForm(myuser, request.POST)
            if form.is_valid():
                form.save()
                return redirect('password_reset_complete')
        return render(request, 'changepassword.html', {'form': form})

    else:
        return HttpResponse('Invalid token')
